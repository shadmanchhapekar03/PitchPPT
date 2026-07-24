from crewai.tools import BaseTool
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
import json, os, requests, urllib.parse


class pitchDeck(BaseTool):
    name: str = "PPTPitch Generation Tool"
    description: str = (
        "Generates a PowerPoint (.pptx) pitch deck file from a JSON string. "
        "Input MUST be a valid JSON string with this exact structure:\n"
        '{"title": "Deck Title", "subtitle": "Deck Subtitle", '
        '"theme": {"background_color": "1A1A2E", "accent_color": "E94E77", "text_color": "FFFFFF"}, '
        '"slides": [{"title": "Slide Title", "content": ["bullet 1", "bullet 2"], '
        '"image_prompt": "short visual description"}]}\n'
        "Returns the file path of the generated .pptx file."
    )

    def _generate_image(self, prompt: str, save_path: str, width=800, height=800) -> str | None:
        try:
            encoded = urllib.parse.quote(prompt)
            url = f"https://image.pollinations.ai/prompt/{encoded}?width={width}&height={height}&nologo=true"
            resp = requests.get(url, timeout=60)
            resp.raise_for_status()
            with open(save_path, "wb") as f:
                f.write(resp.content)
            return save_path
        except Exception as e:
            print(f"Image generation failed: {e}")
            return None

    def _style_background(self, slide, color: RGBColor):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _add_accent_bar(self, slide, prs, color: RGBColor):
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(150000), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    def _style_text_frame(self, text_frame, color: RGBColor, size=18, bold=False):
        for p in text_frame.paragraphs:
            p.font.color.rgb = color
            p.font.size = Pt(size)
            p.font.bold = bold

    def _run(self, ppt_json: str) -> str:
        try:
            data = json.loads(ppt_json)
        except json.JSONDecodeError as e:
            return f"Invalid JSON format: {e}"

        theme = data.get("theme", {})
        BG_COLOR = RGBColor.from_string(theme.get("background_color", "1A1A2E"))
        ACCENT_COLOR = RGBColor.from_string(theme.get("accent_color", "E94E77"))
        TEXT_COLOR = RGBColor.from_string(theme.get("text_color", "FFFFFF"))

        os.makedirs("generated_files", exist_ok=True)
        safe_name = "".join(c if c.isalnum() else "_" for c in data.get("title", "deck"))[:40]
        img_dir = f"generated_files/{safe_name}_images"
        os.makedirs(img_dir, exist_ok=True)

        prs = Presentation()

        # --- Title Slide ---
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        self._style_background(slide, BG_COLOR)
        self._add_accent_bar(slide, prs, ACCENT_COLOR)

        slide.shapes.title.text = data.get('title', 'Untitled Deck')
        self._style_text_frame(slide.shapes.title.text_frame, TEXT_COLOR, size=36, bold=True)

        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = data.get('subtitle', '')
            self._style_text_frame(slide.placeholders[1].text_frame, ACCENT_COLOR, size=18)

        # --- Content Slides ---
        blank_layout = prs.slide_layouts[6]  # blank layout, full manual control
        SLIDE_W = prs.slide_width
        SLIDE_H = prs.slide_height

        for idx, s in enumerate(data.get("slides", [])):
            slide = prs.slides.add_slide(blank_layout)
            self._style_background(slide, BG_COLOR)
            self._add_accent_bar(slide, prs, ACCENT_COLOR)

            # --- Title (top, full width, fixed height so nothing overlaps it) ---
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9.0), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.text = s.get("title", "")
            self._style_text_frame(tf, TEXT_COLOR, size=28, bold=True)

            # --- Image (left half, BELOW title) ---
            image_prompt = s.get("image_prompt", s.get("title", "business concept"))
            img_dir_local = img_dir
            img_path = f"{img_dir_local}/slide_{idx}.png"
            result_path = self._generate_image(image_prompt, img_path)
            if result_path:
                try:
                    slide.shapes.add_picture(result_path, Inches(0.5), Inches(1.4), width=Inches(4.0))
                except Exception as e:
                    print(f"Could not insert image on slide {idx}: {e}")

            # --- Content (right half, BELOW title, independent box) ---
            content_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.4), Inches(4.9), Inches(5.3))
            body = content_box.text_frame
            body.word_wrap = True
            content = s.get("content", [])
            for i, point in enumerate(content):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = "•  " + point
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)

        # --- Save File ---
        file_path = f"generated_files/{safe_name}.pptx"
        prs.save(file_path)
        print(file_path)

        return f"Presentation successfully created at: {file_path}"